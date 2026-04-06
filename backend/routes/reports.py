"""
Report generation routes for Causal Investigations.
Generates PowerPoint and PDF reports.
"""
import os
from fastapi import APIRouter, Depends, HTTPException, Response
from pydantic import BaseModel
from typing import Optional
from datetime import datetime, timezone
import io
import logging
from database import db
from auth import get_current_user

# AI Integration
from openai import OpenAI

# PowerPoint imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# PDF imports
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT

logger = logging.getLogger(__name__)
router = APIRouter(tags=["Reports"])

# Color scheme
COLORS = {
    "primary": RGBColor(79, 70, 229),  # Indigo-600
    "secondary": RGBColor(100, 116, 139),  # Slate-500
    "danger": RGBColor(239, 68, 68),  # Red-500
    "warning": RGBColor(245, 158, 11),  # Amber-500
    "success": RGBColor(34, 197, 94),  # Green-500
    "white": RGBColor(255, 255, 255),
    "dark": RGBColor(30, 41, 59),  # Slate-800
}


async def get_investigation_data(investigation_id: str, user_id: str):
    """Fetch all investigation data for report generation."""
    # Get investigation
    investigation = await db.investigations.find_one(
        {"id": investigation_id, "created_by": user_id},
        {"_id": 0}
    )
    if not investigation:
        raise HTTPException(status_code=404, detail="Investigation not found")
    
    # Get timeline events
    events = await db.timeline_events.find(
        {"investigation_id": investigation_id},
        {"_id": 0}
    ).sort("event_time", 1).to_list(100)
    
    # Get failure identifications
    failures = await db.failure_identifications.find(
        {"investigation_id": investigation_id},
        {"_id": 0}
    ).to_list(100)
    
    # Get cause nodes
    causes = await db.cause_nodes.find(
        {"investigation_id": investigation_id},
        {"_id": 0}
    ).to_list(100)
    
    # Get action items
    actions = await db.action_items.find(
        {"investigation_id": investigation_id},
        {"_id": 0}
    ).to_list(100)
    
    return {
        "investigation": investigation,
        "events": events,
        "failures": failures,
        "causes": causes,
        "actions": actions,
    }


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide to the presentation."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["primary"]
    shape.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS["white"]
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
        sub_frame = sub_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = subtitle
        sub_para.font.size = Pt(24)
        sub_para.font.color.rgb = COLORS["white"]
        sub_para.alignment = PP_ALIGN.CENTER
    
    # Add footer watermark
    add_slide_footer(slide, prs, is_title_slide=True)
    
    return slide


def add_slide_footer(slide, prs, is_title_slide=False):
    """Add 'Generated with AssetIQ' footer to a slide."""
    footer_box = slide.shapes.add_textbox(
        Inches(7), Inches(7.1), Inches(2.8), Inches(0.3)
    )
    footer_frame = footer_box.text_frame
    footer_para = footer_frame.paragraphs[0]
    footer_para.text = "Generated with AssetIQ"
    footer_para.font.size = Pt(9)
    footer_para.font.italic = True
    footer_para.alignment = PP_ALIGN.RIGHT
    if is_title_slide:
        footer_para.font.color.rgb = RGBColor(200, 200, 220)  # Light on dark
    else:
        footer_para.font.color.rgb = RGBColor(148, 163, 184)  # slate-400


def add_section_slide(prs, title):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8), Inches(10), Inches(1.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["primary"]
    bar.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS["white"]
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add footer
    add_slide_footer(slide, prs)
    
    return slide


def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["primary"]
    bar.line.fill.background()
    
    # Add title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS["white"]
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items[:8]):  # Limit to 8 items per slide
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS["dark"]
        p.space_after = Pt(12)
    
    # Add footer
    add_slide_footer(slide, prs)
    
    return slide


def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["primary"]
    bar.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS["white"]
    
    # Add table
    num_rows = min(len(rows) + 1, 10)  # Limit rows
    num_cols = len(headers)
    
    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.3), Inches(1.5),
        Inches(9.4), Inches(5)
    ).table
    
    # Style header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["primary"]
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(14)
        para.font.color.rgb = COLORS["white"]
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Fill data rows
    for row_idx, row_data in enumerate(rows[:num_rows-1]):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)[:50] if cell_text else "-"
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = COLORS["dark"]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Add footer
    add_slide_footer(slide, prs)
    
    return slide


def generate_pptx(data):
    """Generate PowerPoint presentation from investigation data."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    inv = data["investigation"]
    events = data["events"]
    failures = data["failures"]
    causes = data["causes"]
    actions = data["actions"]
    
    # Title Slide
    add_title_slide(
        prs,
        inv.get("title", "Investigation Report"),
        f"Case: {inv.get('case_number', 'N/A')} | {inv.get('incident_date', '')[:10] if inv.get('incident_date') else 'Date N/A'}"
    )
    
    # Overview Slide
    overview_items = [
        f"Asset: {inv.get('asset_name', 'N/A')}",
        f"Location: {inv.get('location', 'N/A')}",
        f"Lead Investigator: {inv.get('investigation_leader', 'N/A')}",
        f"Status: {inv.get('status', 'N/A').replace('_', ' ').title()}",
        f"Created: {inv.get('created_at', '')[:10] if inv.get('created_at') else 'N/A'}",
    ]
    if inv.get("description"):
        overview_items.append(f"Description: {inv.get('description')[:100]}...")
    
    add_content_slide(prs, "Investigation Overview", overview_items)
    
    # Timeline Section
    if events:
        add_section_slide(prs, "Timeline of Events")
        
        event_headers = ["Time", "Category", "Description", "Confidence"]
        event_rows = []
        for event in events:
            event_rows.append([
                event.get("event_time", "")[:16] if event.get("event_time") else "-",
                event.get("category", "-").replace("_", " ").title(),
                event.get("description", "-")[:40],
                event.get("confidence", "-").title()
            ])
        
        add_table_slide(prs, f"Events Timeline ({len(events)} events)", event_headers, event_rows)
    
    # Failures Section
    if failures:
        add_section_slide(prs, "Failure Identifications")
        
        failure_headers = ["Asset", "Component", "Failure Mode", "Mechanism"]
        failure_rows = []
        for f in failures:
            failure_rows.append([
                f.get("asset_name", "-")[:20],
                f.get("component", "-")[:20],
                f.get("failure_mode", "-")[:25],
                f.get("degradation_mechanism", "-")[:25]
            ])
        
        add_table_slide(prs, f"Identified Failures ({len(failures)})", failure_headers, failure_rows)
    
    # Root Cause Analysis Section
    if causes:
        add_section_slide(prs, "Root Cause Analysis")
        
        root_causes = [c for c in causes if c.get("is_root_cause")]
        contributing = [c for c in causes if not c.get("is_root_cause")]
        
        if root_causes:
            cause_items = [f"[ROOT] {c.get('description', 'N/A')[:80]}" for c in root_causes]
            add_content_slide(prs, f"Root Causes ({len(root_causes)})", cause_items)
        
        if contributing:
            contrib_items = [f"{c.get('description', 'N/A')[:80]}" for c in contributing[:8]]
            add_content_slide(prs, f"Contributing Factors ({len(contributing)})", contrib_items)
    
    # Actions Section
    if actions:
        add_section_slide(prs, "Corrective Actions")
        
        action_headers = ["#", "Description", "Owner", "Priority", "Status"]
        action_rows = []
        for a in actions:
            action_rows.append([
                a.get("action_number", "-"),
                a.get("description", "-")[:30],
                a.get("owner", "-")[:15],
                a.get("priority", "-").title(),
                a.get("status", "-").replace("_", " ").title()
            ])
        
        add_table_slide(prs, f"Action Items ({len(actions)})", action_headers, action_rows)
    
    # Summary Slide
    summary_items = [
        f"Total Events: {len(events)}",
        f"Identified Failures: {len(failures)}",
        f"Root Causes Found: {len([c for c in causes if c.get('is_root_cause')])}",
        f"Contributing Factors: {len([c for c in causes if not c.get('is_root_cause')])}",
        f"Action Items: {len(actions)}",
        f"Completed Actions: {len([a for a in actions if a.get('status') == 'completed'])}",
    ]
    add_content_slide(prs, "Investigation Summary", summary_items)
    
    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def generate_pdf(data):
    """Generate PDF report from investigation data."""
    inv = data["investigation"]
    events = data["events"]
    failures = data["failures"]
    causes = data["causes"]
    actions = data["actions"]
    
    buffer = io.BytesIO()
    
    # Custom page template with footer
    def add_page_footer(canvas, doc):
        canvas.saveState()
        # Right-aligned footer
        canvas.setFont('Helvetica-Oblique', 8)
        canvas.setFillColor(colors.HexColor('#94A3B8'))
        canvas.drawRightString(
            letter[0] - 72,  # Right margin
            40,  # 40 points from bottom
            "Generated with AssetIQ"
        )
        canvas.restoreState()
    
    doc = SimpleDocTemplate(
        buffer,
        pagesize=letter,
        rightMargin=72,
        leftMargin=72,
        topMargin=72,
        bottomMargin=72
    )
    
    # Styles
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        spaceAfter=20,
        textColor=colors.HexColor('#4F46E5'),
        alignment=TA_CENTER
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=16,
        spaceAfter=12,
        spaceBefore=20,
        textColor=colors.HexColor('#4F46E5')
    )
    
    subheading_style = ParagraphStyle(
        'CustomSubheading',
        parent=styles['Heading3'],
        fontSize=12,
        spaceAfter=8,
        spaceBefore=12,
        textColor=colors.HexColor('#1E293B')
    )
    
    normal_style = ParagraphStyle(
        'CustomNormal',
        parent=styles['Normal'],
        fontSize=10,
        spaceAfter=6,
        textColor=colors.HexColor('#334155')
    )
    
    story = []
    
    # Title
    story.append(Paragraph(inv.get("title", "Investigation Report"), title_style))
    story.append(Paragraph(f"Case Number: {inv.get('case_number', 'N/A')}", normal_style))
    story.append(Spacer(1, 20))
    
    # Overview Table
    story.append(Paragraph("Investigation Overview", heading_style))
    overview_data = [
        ["Field", "Value"],
        ["Asset", inv.get("asset_name", "N/A")],
        ["Location", inv.get("location", "N/A")],
        ["Incident Date", inv.get("incident_date", "N/A")[:10] if inv.get("incident_date") else "N/A"],
        ["Lead Investigator", inv.get("investigation_leader", "N/A")],
        ["Status", inv.get("status", "N/A").replace("_", " ").title()],
    ]
    
    overview_table = Table(overview_data, colWidths=[2*inch, 4*inch])
    overview_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4F46E5')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
        ('TOPPADDING', (0, 0), (-1, -1), 8),
        ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#F8FAFC')),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#CBD5E1')),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
    ]))
    story.append(overview_table)
    
    if inv.get("description"):
        story.append(Spacer(1, 10))
        story.append(Paragraph("Description:", subheading_style))
        story.append(Paragraph(inv.get("description"), normal_style))
    
    # Timeline Events
    if events:
        story.append(PageBreak())
        story.append(Paragraph(f"Timeline of Events ({len(events)})", heading_style))
        
        event_data = [["Time", "Category", "Description", "Confidence"]]
        for event in events:
            event_data.append([
                event.get("event_time", "-")[:16] if event.get("event_time") else "-",
                event.get("category", "-").replace("_", " ").title()[:15],
                event.get("description", "-")[:50],
                event.get("confidence", "-").title()
            ])
        
        event_table = Table(event_data, colWidths=[1.2*inch, 1.2*inch, 2.8*inch, 0.8*inch])
        event_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4F46E5')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 9),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#F8FAFC')),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#CBD5E1')),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
        ]))
        story.append(event_table)
    
    # Failures
    if failures:
        story.append(PageBreak())
        story.append(Paragraph(f"Failure Identifications ({len(failures)})", heading_style))
        
        failure_data = [["Asset", "Component", "Failure Mode", "Mechanism"]]
        for f in failures:
            failure_data.append([
                f.get("asset_name", "-")[:20],
                f.get("component", "-")[:15],
                f.get("failure_mode", "-")[:25],
                f.get("degradation_mechanism", "-")[:25]
            ])
        
        failure_table = Table(failure_data, colWidths=[1.5*inch, 1.2*inch, 1.8*inch, 1.5*inch])
        failure_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4F46E5')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 9),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#F8FAFC')),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#CBD5E1')),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
        ]))
        story.append(failure_table)
    
    # Root Causes
    if causes:
        story.append(PageBreak())
        story.append(Paragraph("Root Cause Analysis", heading_style))
        
        root_causes = [c for c in causes if c.get("is_root_cause")]
        contributing = [c for c in causes if not c.get("is_root_cause")]
        
        if root_causes:
            story.append(Paragraph(f"Root Causes ({len(root_causes)})", subheading_style))
            for c in root_causes:
                story.append(Paragraph(f"• <b>[ROOT]</b> {c.get('description', 'N/A')}", normal_style))
        
        if contributing:
            story.append(Spacer(1, 10))
            story.append(Paragraph(f"Contributing Factors ({len(contributing)})", subheading_style))
            for c in contributing:
                category = c.get("category", "").replace("_", " ").title()
                story.append(Paragraph(f"• [{category}] {c.get('description', 'N/A')}", normal_style))
    
    # Actions
    if actions:
        story.append(PageBreak())
        story.append(Paragraph(f"Corrective Actions ({len(actions)})", heading_style))
        
        action_data = [["#", "Description", "Owner", "Priority", "Due Date", "Status"]]
        for a in actions:
            action_data.append([
                a.get("action_number", "-"),
                a.get("description", "-")[:35],
                a.get("owner", "-")[:12],
                a.get("priority", "-").title(),
                a.get("due_date", "-")[:10] if a.get("due_date") else "-",
                a.get("status", "-").replace("_", " ").title()
            ])
        
        action_table = Table(action_data, colWidths=[0.6*inch, 2.2*inch, 1*inch, 0.8*inch, 0.8*inch, 0.8*inch])
        action_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4F46E5')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 8),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
            ('TOPPADDING', (0, 0), (-1, -1), 5),
            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#F8FAFC')),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#CBD5E1')),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
        ]))
        story.append(action_table)
    
    # Summary
    story.append(PageBreak())
    story.append(Paragraph("Investigation Summary", heading_style))
    
    summary_data = [
        ["Metric", "Count"],
        ["Timeline Events", str(len(events))],
        ["Identified Failures", str(len(failures))],
        ["Root Causes", str(len([c for c in causes if c.get("is_root_cause")]))],
        ["Contributing Factors", str(len([c for c in causes if not c.get("is_root_cause")]))],
        ["Total Actions", str(len(actions))],
        ["Completed Actions", str(len([a for a in actions if a.get("status") == "completed"]))],
    ]
    
    summary_table = Table(summary_data, colWidths=[3*inch, 2*inch])
    summary_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4F46E5')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
        ('TOPPADDING', (0, 0), (-1, -1), 8),
        ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#F8FAFC')),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#CBD5E1')),
        ('ALIGN', (0, 0), (0, -1), 'LEFT'),
        ('ALIGN', (1, 0), (1, -1), 'CENTER'),
    ]))
    story.append(summary_table)
    
    # Footer
    story.append(Spacer(1, 30))
    story.append(Paragraph(
        f"Generated on {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')} | AssetIQ Causal Engine",
        ParagraphStyle('Footer', parent=normal_style, alignment=TA_CENTER, textColor=colors.HexColor('#94A3B8'))
    ))
    
    doc.build(story, onFirstPage=add_page_footer, onLaterPages=add_page_footer)
    buffer.seek(0)
    return buffer.getvalue()


@router.get("/investigations/{investigation_id}/report/pptx")
async def generate_pptx_report(
    investigation_id: str,
    current_user: dict = Depends(get_current_user)
):
    """Generate PowerPoint report for an investigation."""
    try:
        data = await get_investigation_data(investigation_id, current_user["id"])
        pptx_bytes = generate_pptx(data)
        
        filename = f"Investigation_{data['investigation'].get('case_number', investigation_id)}.pptx"
        
        return Response(
            content=pptx_bytes,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error generating PPTX report: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate report: {str(e)}")


@router.get("/investigations/{investigation_id}/report/pdf")
async def generate_pdf_report(
    investigation_id: str,
    current_user: dict = Depends(get_current_user)
):
    """Generate PDF report for an investigation."""
    try:
        data = await get_investigation_data(investigation_id, current_user["id"])
        pdf_bytes = generate_pdf(data)
        
        filename = f"Investigation_{data['investigation'].get('case_number', investigation_id)}.pdf"
        
        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error generating PDF report: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate report: {str(e)}")


# ============= AI SUMMARY =============

class AISummaryResponse(BaseModel):
    summary: str
    key_findings: list
    next_steps: list
    recommendations: list


async def generate_ai_summary(data: dict) -> dict:
    """Generate AI-powered summary of investigation with next steps."""
    inv = data["investigation"]
    events = data["events"]
    failures = data["failures"]
    causes = data["causes"]
    actions = data["actions"]
    
    # Compute statistics for richer context
    root_causes = [c for c in causes if c.get('is_root_cause')]
    contributing_factors = [c for c in causes if not c.get('is_root_cause')]
    open_actions = [a for a in actions if a.get('status') in ['open', 'in_progress']]
    completed_actions = [a for a in actions if a.get('status') == 'completed']
    high_priority_actions = [a for a in actions if a.get('priority') == 'high']
    validated_actions = [a for a in actions if a.get('is_validated')]
    
    # Build detailed timeline narrative
    timeline_details = []
    for i, e in enumerate(events[:15], 1):
        event_time = e.get('event_time', '')[:16] if e.get('event_time') else 'Unknown time'
        category = e.get('category', 'event').replace('_', ' ').title()
        confidence = e.get('confidence', 'medium').title()
        desc = e.get('description', 'No description')
        timeline_details.append(f"{i}. [{event_time}] ({category} - {confidence} confidence): {desc}")
    
    # Build detailed failure analysis
    failure_details = []
    for i, f in enumerate(failures[:10], 1):
        asset = f.get('asset_name', 'Unknown Asset')
        component = f.get('component', 'Unknown Component')
        mode = f.get('failure_mode', 'Unknown Mode')
        mechanism = f.get('degradation_mechanism', 'Unknown Mechanism')
        discipline = f.get('discipline', 'N/A')
        failure_details.append(
            f"{i}. Asset: {asset} | Component: {component}\n"
            f"   - Failure Mode: {mode}\n"
            f"   - Degradation Mechanism: {mechanism}\n"
            f"   - Discipline: {discipline}"
        )
    
    # Build detailed root cause analysis
    root_cause_details = []
    for i, c in enumerate(root_causes[:8], 1):
        category = c.get('category', 'unknown').replace('_', ' ').title()
        desc = c.get('description', 'No description')
        root_cause_details.append(f"{i}. [{category}] {desc}")
    
    # Build contributing factors with categories
    factor_details = []
    factor_categories = {}
    for c in contributing_factors[:12]:
        cat = c.get('category', 'other').replace('_', ' ').title()
        if cat not in factor_categories:
            factor_categories[cat] = []
        factor_categories[cat].append(c.get('description', 'No description'))
    
    for cat, items in factor_categories.items():
        factor_details.append(f"  {cat}:")
        for item in items[:3]:
            factor_details.append(f"    • {item}")
    
    # Build detailed action plan
    action_details = []
    for i, a in enumerate(actions[:12], 1):
        num = a.get('action_number', f'A{i}')
        desc = a.get('description', 'No description')
        owner = a.get('owner', 'Unassigned')
        priority = a.get('priority', 'medium').title()
        status = a.get('status', 'open').replace('_', ' ').title()
        due_date = a.get('due_date', '')[:10] if a.get('due_date') else 'No due date'
        validated = "✓ Validated" if a.get('is_validated') else "Not validated"
        action_details.append(
            f"{num}. {desc}\n"
            f"   - Owner: {owner} | Priority: {priority} | Status: {status}\n"
            f"   - Due: {due_date} | {validated}"
        )
    
    # Prepare comprehensive context for AI
    context = f"""
================================================================================
                        CAUSAL INVESTIGATION ANALYSIS
================================================================================

INVESTIGATION IDENTIFICATION
----------------------------
• Title: {inv.get('title', 'N/A')}
• Case Number: {inv.get('case_number', 'N/A')}
• Asset: {inv.get('asset_name', 'N/A')}
• Location: {inv.get('location', 'N/A')}
• Incident Date: {inv.get('incident_date', 'N/A')[:10] if inv.get('incident_date') else 'N/A'}
• Investigation Status: {inv.get('status', 'N/A').replace('_', ' ').title()}
• Lead Investigator: {inv.get('investigation_leader', 'N/A')}

INCIDENT DESCRIPTION
--------------------
{inv.get('description', 'No description provided.')}

================================================================================
                            TIMELINE OF EVENTS
                          ({len(events)} events recorded)
================================================================================
{chr(10).join(timeline_details) if timeline_details else 'No timeline events recorded.'}

================================================================================
                         FAILURE IDENTIFICATION
                        ({len(failures)} failures identified)
================================================================================
{chr(10).join(failure_details) if failure_details else 'No failures identified.'}

================================================================================
                          ROOT CAUSE ANALYSIS
================================================================================
PRIMARY ROOT CAUSES ({len(root_causes)} identified):
{chr(10).join(root_cause_details) if root_cause_details else '• No root causes identified yet.'}

CONTRIBUTING FACTORS ({len(contributing_factors)} identified):
{chr(10).join(factor_details) if factor_details else '• No contributing factors identified.'}

================================================================================
                        CORRECTIVE ACTION PLAN
                        ({len(actions)} total actions)
================================================================================
Action Statistics:
• Open/In-Progress: {len(open_actions)}
• Completed: {len(completed_actions)}
• High Priority: {len(high_priority_actions)}
• Validated: {len(validated_actions)}

Detailed Action Items:
{chr(10).join(action_details) if action_details else 'No corrective actions defined.'}

================================================================================
"""
    
    prompt = f"""You are a senior reliability engineer and root cause analysis expert preparing an executive briefing based on a completed causal investigation. Your analysis must be DATA-DRIVEN and reference SPECIFIC details from the investigation.

ANALYSIS REQUIREMENTS:
1. EXECUTIVE SUMMARY (3-4 detailed paragraphs):
   - Paragraph 1: Describe what happened - reference the specific asset ({inv.get('asset_name', 'the asset')}), location, and incident date. Summarize the event sequence.
   - Paragraph 2: Explain WHY it happened - explicitly name the identified root causes and explain their relationship to the failure modes.
   - Paragraph 3: Describe the impact and current action plan status. Reference specific numbers (e.g., "{len(open_actions)} actions pending, {len(completed_actions)} completed").
   - Paragraph 4: Provide overall assessment of investigation completeness and risk exposure.

2. KEY FINDINGS (5-7 findings):
   - Each finding must reference specific data from the investigation
   - Include the actual failure modes, mechanisms, and root causes identified
   - Cite specific contributing factors discovered
   - Reference action plan progress with real numbers

3. NEXT STEPS (5-8 specific actions):
   - Prioritize based on the HIGH PRIORITY actions that are still OPEN
   - Reference specific owners and due dates where available
   - Include validation requirements for unvalidated actions
   - Add follow-up investigation steps if root causes are incomplete

4. STRATEGIC RECOMMENDATIONS (3-5 recommendations):
   - Based on the specific failure modes and mechanisms found
   - Reference the contributing factor categories
   - Include preventive measures tied to the actual root causes
   - Suggest systemic improvements based on investigation findings

{context}

CRITICAL: Your response must be CONTENT-RICH and reference actual data points, names, dates, equipment, and findings from this investigation. Avoid generic statements. Every bullet point should contain specific information.

Respond in JSON format:
{{
  "summary": "Multi-paragraph executive summary with specific references...",
  "key_findings": ["Specific finding 1 with data...", "Specific finding 2 with data...", ...],
  "next_steps": ["Specific action with owner/date...", "Specific action 2...", ...],
  "recommendations": ["Specific recommendation based on findings...", ...]
}}"""

    try:
        import uuid
        inv_id = inv.get('id', str(uuid.uuid4()))
        
        api_key = os.environ.get("OPENAI_API_KEY")
        if not api_key:
            raise Exception("OpenAI API key not configured")
        
        client = OpenAI(api_key=api_key)
        
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a senior reliability engineer and root cause analysis expert. Always respond with valid JSON only."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.5
        )
        
        response_text = response.choices[0].message.content.strip()
        
        # Parse JSON response
        import json
        # Clean up response if needed
        if response_text.startswith("```json"):
            response_text = response_text[7:]
        if response_text.startswith("```"):
            response_text = response_text[3:]
        if response_text.endswith("```"):
            response_text = response_text[:-3]
        
        result = json.loads(response_text.strip())
        return result
    except Exception as e:
        logger.error(f"Error generating AI summary: {e}")
        # Return a detailed fallback that still references specific data
        root_cause_list = [c.get('description', 'Unknown') for c in root_causes[:3]]
        failure_list = [f"{f.get('asset_name', 'Asset')}: {f.get('failure_mode', 'Unknown failure')}" for f in failures[:3]]
        pending_owners = list(set([a.get('owner', 'Unassigned') for a in open_actions[:5]]))
        
        return {
            "summary": f"""Investigation '{inv.get('title', 'N/A')}' (Case #{inv.get('case_number', 'N/A')}) was initiated following an incident at {inv.get('asset_name', 'the asset')} located in {inv.get('location', 'the facility')}. The investigation timeline captured {len(events)} events leading to the incident.

The root cause analysis identified {len(root_causes)} primary root cause(s) and {len(contributing_factors)} contributing factors. {f"Key root causes include: {'; '.join(root_cause_list)}." if root_cause_list else "Root cause identification is still in progress."} Failure analysis revealed {len(failures)} failure mode(s) affecting the asset.

The corrective action plan contains {len(actions)} total actions, of which {len(completed_actions)} are completed, {len(open_actions)} remain open, and {len(validated_actions)} have been validated. {f"High priority actions ({len(high_priority_actions)}) require immediate attention." if high_priority_actions else ""}

Overall investigation status: {inv.get('status', 'In Progress').replace('_', ' ').title()}. {"Immediate action required to address open high-priority items." if high_priority_actions else "Continue systematic closure of remaining action items."}""",
            "key_findings": [
                f"Timeline Analysis: {len(events)} events documented, establishing the incident sequence",
                f"Failure Modes: {len(failures)} failure(s) identified - {', '.join(failure_list) if failure_list else 'Analysis in progress'}",
                f"Root Causes: {len(root_causes)} primary root cause(s) confirmed through investigation",
                f"Contributing Factors: {len(contributing_factors)} factors identified across {len(factor_categories)} categories",
                f"Action Plan Status: {len(completed_actions)}/{len(actions)} actions completed ({int(len(completed_actions)/len(actions)*100) if actions else 0}% completion rate)",
                f"Validation Progress: {len(validated_actions)} of {len(actions)} actions validated by designated personnel"
            ],
            "next_steps": [
                f"Complete root cause validation for {len(root_causes)} identified causes" if root_causes else "Initiate root cause analysis to identify primary failure causes",
                f"Address {len(high_priority_actions)} high-priority open actions immediately" if high_priority_actions else "Review action prioritization and ensure critical items are flagged",
                f"Obtain validation sign-off for {len(actions) - len(validated_actions)} unvalidated actions" if actions else "Define corrective actions based on root cause findings",
                f"Follow up with action owners: {', '.join(pending_owners[:3])}" if pending_owners else "Assign owners to all open action items",
                "Schedule 2-week progress review with investigation team",
                "Update FMEA library with lessons learned from this investigation"
            ],
            "recommendations": [
                f"Implement preventive measures targeting the identified root causes: {root_cause_list[0] if root_cause_list else 'TBD'}",
                f"Update inspection protocols for {inv.get('asset_name', 'affected equipment')} based on failure mode analysis",
                "Establish leading indicators to detect similar failure patterns earlier",
                "Share investigation findings with maintenance and operations teams during toolbox talks"
            ]
        }


@router.get("/investigations/{investigation_id}/ai-summary")
async def get_ai_summary(
    investigation_id: str,
    current_user: dict = Depends(get_current_user)
):
    """Generate AI-powered summary of investigation with next steps."""
    try:
        data = await get_investigation_data(investigation_id, current_user["id"])
        summary = await generate_ai_summary(data)
        return summary
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error getting AI summary: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate AI summary: {str(e)}")
